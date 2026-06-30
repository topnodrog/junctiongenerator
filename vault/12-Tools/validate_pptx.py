#!/usr/bin/env python3
from pptx import Presentation

prs = Presentation("JGC_Funding_Pitch.pptx")

print(f"✓ Presentation loaded: {len(prs.slides)} slides\n")

for i, slide in enumerate(prs.slides, 1):
    print(f"## Slide {i}")
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text = shape.text.strip()
            if text:
                # Show first 100 chars
                display = text[:100] + ("..." if len(text) > 100 else "")
                print(f"  {display}")
    print()
